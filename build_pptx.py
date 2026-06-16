"""
Genera la presentacion PowerPoint del proyecto "Casa de la Cultura de Filadelfia",
incluyendo una diapositiva con el codigo QR de acceso al museo virtual.

Uso:
    python build_pptx.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

CARPETA = Path(__file__).parent
IMG = CARPETA / "imagenes"
QR = CARPETA / "qr-museo.png"
SALIDA = CARPETA / "Presentacion-Casa-Cultura-Filadelfia.pptx"

# --- Paleta (misma que el museo virtual en HTML) ---
TERRACOTA = RGBColor(0xA8, 0x5C, 0x32)
TERRACOTA_OSC = RGBColor(0x7C, 0x43, 0x22)
CREMA = RGBColor(0xF4, 0xEC, 0xDF)
CREMA_OSC = RGBColor(0xE8, 0xDC, 0xC4)
CAFE = RGBColor(0x4A, 0x33, 0x24)
CAFE_CLARO = RGBColor(0x6B, 0x4F, 0x3A)
VERDE = RGBColor(0x5C, 0x6B, 0x3F)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)

FUENTE_TITULO = "Trebuchet MS"
FUENTE_TEXTO = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def fondo(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def caja_texto(slide, izq, arr, ancho, alto, anclaje=MSO_ANCHOR.TOP):
    cuadro = slide.shapes.add_textbox(izq, arr, ancho, alto)
    cuadro.text_frame.vertical_anchor = anclaje
    cuadro.text_frame.word_wrap = True
    return cuadro


def parrafo(frame, texto, tam=18, color=CAFE, negrita=False, fuente=FUENTE_TEXTO,
            alineacion=PP_ALIGN.LEFT, espacio_antes=0, primero=False, cursiva=False):
    p = frame.paragraphs[0] if primero and not frame.paragraphs[0].runs else frame.add_paragraph()
    p.alignment = alineacion
    p.space_before = Pt(espacio_antes)
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(tam)
    run.font.bold = negrita
    run.font.italic = cursiva
    run.font.name = fuente
    run.font.color.rgb = color
    return p


def franja_lateral(slide, color=TERRACOTA, ancho=Inches(0.28)):
    franja = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, ancho, prs.slide_height)
    franja.fill.solid()
    franja.fill.fore_color.rgb = color
    franja.line.fill.background()
    franja.shadow.inherit = False


def encabezado(slide, titulo, subtitulo=None):
    franja_lateral(slide)
    caja = caja_texto(slide, Inches(0.75), Inches(0.45), Inches(11.8), Inches(1.2))
    parrafo(caja.text_frame, titulo, tam=32, color=TERRACOTA_OSC, negrita=True,
            fuente=FUENTE_TITULO, primero=True)
    if subtitulo:
        caja2 = caja_texto(slide, Inches(0.75), Inches(1.15), Inches(11.8), Inches(0.6))
        parrafo(caja2.text_frame, subtitulo, tam=16, color=CAFE_CLARO, cursiva=True, primero=True)


def pie_pagina(slide, texto):
    caja = caja_texto(slide, Inches(0.75), Inches(7.05), Inches(11.8), Inches(0.35))
    parrafo(caja.text_frame, texto, tam=10.5, color=CAFE_CLARO, primero=True)


def viñetas(slide, items, izq=Inches(0.85), arr=Inches(1.9), ancho=Inches(11.4), alto=Inches(4.8),
            tam=18, espacio=10):
    caja = caja_texto(slide, izq, arr, ancho, alto)
    frame = caja.text_frame
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            titulo_item, detalle = item
        else:
            titulo_item, detalle = item, None
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.space_before = Pt(espacio if i else 0)
        run = p.add_run()
        run.text = f"●  {titulo_item}"
        run.font.size = Pt(tam)
        run.font.bold = True
        run.font.name = FUENTE_TEXTO
        run.font.color.rgb = TERRACOTA_OSC
        if detalle:
            p2 = frame.add_paragraph()
            p2.space_before = Pt(2)
            p2.level = 1
            run2 = p2.add_run()
            run2.text = detalle
            run2.font.size = Pt(tam - 3)
            run2.font.name = FUENTE_TEXTO
            run2.font.color.rgb = CAFE
    return caja


def imagen_recortada(slide, ruta, izq, arr, ancho, alto):
    """Inserta una imagen llenando el rectangulo dado (recorte tipo 'cover')."""
    from PIL import Image
    with Image.open(ruta) as im:
        w_px, h_px = im.size
    rel_caja = ancho / alto
    rel_img = w_px / h_px
    pic = slide.shapes.add_picture(str(ruta), izq, arr, height=alto)
    if rel_img < rel_caja:
        pic.height = alto
        pic.width = Emu(int(alto * rel_img))
        nuevo_ancho = Emu(int(alto * rel_img))
        recorte = (nuevo_ancho - ancho) / 2 / nuevo_ancho
        pic.crop_left = recorte
        pic.crop_right = recorte
        pic.width = ancho
        pic.left = izq
    else:
        pic.width = ancho
        pic.height = Emu(int(ancho / rel_img))
        nuevo_alto = Emu(int(ancho / rel_img))
        recorte = (nuevo_alto - alto) / 2 / nuevo_alto
        pic.crop_top = recorte
        pic.crop_bottom = recorte
        pic.height = alto
        pic.top = arr
    pic.left, pic.top = izq, arr
    return pic


# ============================================================
# 1. Portada
# ============================================================
s = prs.slides.add_slide(BLANK)
fondo(s, CAFE)
imagen_recortada(s, IMG / "foto-20.jpg", 0, 0, prs.slide_width, prs.slide_height)

velo = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
velo.fill.solid()
velo.fill.fore_color.rgb = RGBColor(0x28, 0x1A, 0x10)
velo.fill.transparency = 0
velo.line.fill.background()
velo.shadow.inherit = False
velo.fill.fore_color.rgb = RGBColor(0x28, 0x1A, 0x10)
sp = velo.fill._xPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
from pptx.oxml.ns import qn
alpha = sp.makeelement(qn('a:alpha'), {'val': '68000'})
sp.append(alpha)

caja = caja_texto(s, Inches(1), Inches(2.5), Inches(11.3), Inches(2.6), MSO_ANCHOR.MIDDLE)
parrafo(caja.text_frame, "UNIVERSIDAD DE COSTA RICA", tam=15, color=CREMA, fuente=FUENTE_TITULO,
        alineacion=PP_ALIGN.CENTER, primero=True)
parrafo(caja.text_frame, "SR-0033 Seminario de Realidad Nacional II · Patrimonio Cultural", tam=13,
        color=CREMA, fuente=FUENTE_TITULO, alineacion=PP_ALIGN.CENTER, espacio_antes=4)
parrafo(caja.text_frame, "Casa de la Cultura de Filadelfia", tam=44, color=BLANCO, negrita=True,
        fuente=FUENTE_TITULO, alineacion=PP_ALIGN.CENTER, espacio_antes=18)
parrafo(caja.text_frame, "Patrimonio histórico-arquitectónico de Guanacaste, Carrillo · 2026", tam=18,
        color=CREMA_OSC, cursiva=True, alineacion=PP_ALIGN.CENTER, espacio_antes=10)

caja2 = caja_texto(s, Inches(1), Inches(6.3), Inches(11.3), Inches(0.9), MSO_ANCHOR.MIDDLE)
parrafo(caja2.text_frame, "Keiry Campos Castro · Katherine Polanco Landeros · Karen Picado Espinoza · Kendall Quirós Alvarez",
        tam=14, color=CREMA, alineacion=PP_ALIGN.CENTER, primero=True)

# ============================================================
# 2. Justificación
# ============================================================
s = prs.slides.add_slide(BLANK)
fondo(s, CREMA)
encabezado(s, "Justificación", "¿Por qué investigar la Casa de la Cultura de Filadelfia?")
viñetas(s, [
    ("Patrimonio protegido por ley", "Declarada Patrimonio Histórico-Arquitectónico de Costa Rica desde el año 2000 (Decreto Ejecutivo N.° 28698-C)."),
    ("Construcción tradicional en riesgo", "Edificada a mediados del siglo XIX con bahareque y madera, materiales típicos costarricenses hoy deteriorados."),
    ("Deterioro y abandono actual", "El edificio pertenece a la Municipalidad de Carrillo, pero se encuentra sin uso cultural y en mal estado."),
    ("Pérdida de memoria e identidad", "El abandono de un patrimonio así no solo afecta la construcción, sino la historia y memoria de Filadelfia."),
], tam=19)
pie_pagina(s, "Vidrios y Celosías — no aplica · Casa de la Cultura de Filadelfia, 2026")

# ============================================================
# 3. Objetivos
# ============================================================
s = prs.slides.add_slide(BLANK)
fondo(s, CREMA)
encabezado(s, "Objetivos")

caja = caja_texto(s, Inches(0.85), Inches(1.85), Inches(11.4), Inches(0.5))
parrafo(caja.text_frame, "General", tam=20, color=VERDE, negrita=True, fuente=FUENTE_TITULO, primero=True)
caja2 = caja_texto(s, Inches(0.85), Inches(2.3), Inches(11.4), Inches(1.0))
parrafo(caja2.text_frame,
        "Analizar la situación actual de la Casa de la Cultura de Filadelfia como patrimonio "
        "histórico-arquitectónico, con el fin de visibilizar su importancia cultural y el riesgo en el que se encuentra.",
        tam=18, primero=True)

caja3 = caja_texto(s, Inches(0.85), Inches(3.55), Inches(11.4), Inches(0.5))
parrafo(caja3.text_frame, "Específicos", tam=20, color=VERDE, negrita=True, fuente=FUENTE_TITULO, primero=True)
viñetas(s, [
    "Investigar las características históricas y arquitectónicas de la Casa de la Cultura y su declaratoria como patrimonio nacional.",
    "Identificar el estado actual de conservación del edificio y las razones de su deterioro y falta de uso cultural.",
    "Reflexionar sobre el impacto que el abandono de este edificio tiene en la identidad y memoria histórica de la comunidad.",
], izq=Inches(0.85), arr=Inches(4.05), alto=Inches(2.7), tam=17, espacio=8)
pie_pagina(s, "Casa de la Cultura de Filadelfia, 2026")

# ============================================================
# 4. Metodología
# ============================================================
s = prs.slides.add_slide(BLANK)
fondo(s, CREMA)
encabezado(s, "Metodología", "Enfoque cualitativo · Estudio de caso · Diseño no experimental y transversal")
viñetas(s, [
    ("Población y muestra", "Residentes de Filadelfia (10–30 personas) y funcionarios de la Municipalidad de Carrillo (3–5 personas), mediante muestreo intencional."),
    ("Entrevistas semiestructuradas", "25–40 minutos, organizadas en bloques de datos generales y preguntas abiertas sobre memoria, conservación y propuestas."),
    ("Observación directa y sistemática", "Visitas de campo para documentar el estado físico del inmueble (fotografías de la Figura 1 a la 20 del informe)."),
    ("Consideraciones éticas", "Consentimiento informado, voluntariedad, confidencialidad y los principios de autonomía, justicia, beneficencia y no maleficencia."),
], tam=18)
pie_pagina(s, "Casa de la Cultura de Filadelfia, 2026")

# ============================================================
# 5. Estado actual de conservación (con imagen)
# ============================================================
s = prs.slides.add_slide(BLANK)
fondo(s, CREMA)
encabezado(s, "Estado actual de conservación", "Hallazgos de las visitas de campo")
viñetas(s, [
    ("Deterioro del bahareque", "Desprendimientos que dejan sectores expuestos a la humedad y a las variaciones climáticas."),
    ("Daños en madera", "Puertas y escaleras del segundo piso con desgaste que compromete su funcionalidad y seguridad."),
    ("Falta de mantenimiento", "Ausencia de labores periódicas de limpieza, restauración y seguimiento."),
    ("Sin uso cultural activo", "El segundo piso se usó para ensayos de baile hasta cerca del año 2020; hoy permanece inactivo."),
], izq=Inches(0.85), arr=Inches(1.85), ancho=Inches(6.6), alto=Inches(4.9), tam=16.5, espacio=12)
imagen_recortada(s, IMG / "foto-10.jpg", Inches(8.0), Inches(1.85), Inches(4.55), Inches(4.9))
pie_pagina(s, "Casa de la Cultura de Filadelfia, 2026")

# ============================================================
# 6. Memoria viva de la comunidad (con imagen + cita)
# ============================================================
s = prs.slides.add_slide(BLANK)
fondo(s, CREMA)
encabezado(s, "Memoria viva de la comunidad", "Identidad cultural, recuerdos y preocupación por el abandono")
imagen_recortada(s, IMG / "foto-41.jpg", Inches(0.85), Inches(1.85), Inches(4.6), Inches(4.9))

caja = caja_texto(s, Inches(5.75), Inches(1.95), Inches(6.8), Inches(4.7))
frame = caja.text_frame
frame.word_wrap = True
p = frame.paragraphs[0]
run = p.add_run()
run.text = ('"Durante aproximadamente veinte años me dediqué al cuidado de la casa y a la formación de '
            'grupos de baile folclórico, transmitiendo bailes como ‘La Marcha Esquivel’, una manifestación '
            'cultural poco conocida hoy."')
run.font.italic = True
run.font.size = Pt(19)
run.font.name = FUENTE_TEXTO
run.font.color.rgb = CAFE
p2 = frame.add_paragraph()
p2.space_before = Pt(8)
r2 = p2.add_run()
r2.text = "— Doña Magdalena, vecina de Filadelfia"
r2.font.bold = True
r2.font.size = Pt(15)
r2.font.name = FUENTE_TEXTO
r2.font.color.rgb = TERRACOTA_OSC

p3 = frame.add_paragraph()
p3.space_before = Pt(22)
r3 = p3.add_run()
r3.text = ("Las entrevistas revelaron tres ejes centrales: la identidad cultural que representa el edificio, "
           "los recuerdos y experiencias vividos en él, y la preocupación de la comunidad por su abandono y "
           "la pérdida de actividades que antes conectaban a las nuevas generaciones con su historia.")
r3.font.size = Pt(16.5)
r3.font.name = FUENTE_TEXTO
r3.font.color.rgb = CAFE
pie_pagina(s, "Casa de la Cultura de Filadelfia, 2026")

# ============================================================
# 7. Propuestas de la comunidad
# ============================================================
s = prs.slides.add_slide(BLANK)
fondo(s, CREMA)
encabezado(s, "Propuestas de la comunidad", "Todavía se puede recuperar este patrimonio")
viñetas(s, [
    ("Restauración del edificio", "Reparar el bahareque, la madera y la estructura antes de que el deterioro avance más."),
    ("Reapertura para actividades culturales", "Recuperar bailes folclóricos, música tradicional y reuniones comunales."),
    ("Creación de un museo local", "Exhibir fotografías, documentos y objetos históricos de Filadelfia."),
    ("Programas educativos", "Talleres y charlas para que niños y jóvenes conozcan su historia."),
    ("Trabajo conjunto", "Coordinación entre la Municipalidad de Carrillo, el Ministerio de Cultura y la comunidad."),
], tam=18, espacio=11)
pie_pagina(s, "Casa de la Cultura de Filadelfia, 2026")

# ============================================================
# 8. Museo virtual + QR
# ============================================================
s = prs.slides.add_slide(BLANK)
fondo(s, CAFE)
franja_lateral(s, TERRACOTA)

caja = caja_texto(s, Inches(0.85), Inches(0.6), Inches(7.3), Inches(2.6))
parrafo(caja.text_frame, "Museo Virtual", tam=38, color=BLANCO, negrita=True, fuente=FUENTE_TITULO, primero=True)
parrafo(caja.text_frame, "Casa de la Cultura de Filadelfia", tam=20, color=CREMA_OSC, cursiva=True, espacio_antes=4)
parrafo(caja.text_frame,
        "Escanea el código QR con la cámara de tu celular para recorrer la galería completa de "
        "fotografías y conocer la historia, el estado actual y las propuestas de la comunidad.",
        tam=16.5, color=CREMA, espacio_antes=20)
parrafo(caja.text_frame, "Funciona desde cualquier celular con internet, en cualquier momento.",
        tam=14, color=RGBColor(0xE0, 0xB0, 0x80), cursiva=True, espacio_antes=18)
parrafo(caja.text_frame, "kendallquirosalvarez-prog.github.io/museo-casa-cultura-filadelfia",
        tam=12.5, color=CREMA_OSC, espacio_antes=10)

if QR.exists():
    qr_alto = Inches(4.7)
    qr_izq = prs.slide_width - qr_alto - Inches(0.7)
    qr_arr = Emu(int((prs.slide_height - qr_alto) / 2))
    s.shapes.add_picture(str(QR), qr_izq, qr_arr, height=qr_alto)

pie_pagina(s, "Casa de la Cultura de Filadelfia, 2026")
for shp in s.shapes:
    if shp.has_text_frame and "Casa de la Cultura de Filadelfia, 2026" in shp.text_frame.text:
        for p in shp.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = CREMA_OSC

# ============================================================
# 9. Cierre / créditos
# ============================================================
s = prs.slides.add_slide(BLANK)
fondo(s, CREMA)
franja_lateral(s)
caja = caja_texto(s, Inches(0.85), Inches(2.6), Inches(11.4), Inches(2.5), MSO_ANCHOR.MIDDLE)
parrafo(caja.text_frame, "¡Gracias!", tam=40, color=TERRACOTA_OSC, negrita=True, fuente=FUENTE_TITULO,
        alineacion=PP_ALIGN.CENTER, primero=True)
parrafo(caja.text_frame, "Keiry Campos Castro · Katherine Polanco Landeros · Karen Picado Espinoza · Kendall Quirós Alvarez",
        tam=16, color=CAFE, alineacion=PP_ALIGN.CENTER, espacio_antes=16)
parrafo(caja.text_frame, "Universidad de Costa Rica · SR-0033 Seminario de Realidad Nacional II — Patrimonio Cultural · 2026",
        tam=13, color=CAFE_CLARO, cursiva=True, alineacion=PP_ALIGN.CENTER, espacio_antes=6)

prs.save(SALIDA)
print(f"Presentacion guardada en: {SALIDA}")
